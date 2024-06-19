import datetime
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import os

#################################
## 파워포인트 객체 선언
#################################
today = datetime.datetime.today().strftime('%Y%m%d')
prs = Presentation()

#################################
## 제목 슬라이드 추가
#################################
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

# 제목
title = slide.shapes.title
title.text = "주식 보고서"

# 부제목
subtitle = slide.placeholders[1]
subtitle.text = "보고서 작성일 : {date}".format(date=today)

#################################
## 차트 및 테이블 슬라이드 추가
#################################
company = "삼성전자"
df = pd.DataFrame({'close': [50000]})
chart_fname = "D:/Joo.min.hee/IOT강의/Source.1/res/stock_report/삼성전자_chart.png"
table_fname = "D:/Joo.min.hee/IOT강의/Source.1/res/stock_report/삼성전자_table.png"

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

shapes = slide.shapes
shapes.title.text = '{company}, {close}원에 거래 마감'.format(company=company, close=df.iloc[0]['close'])

# 차트 추가
left = Inches(0.5)
height = Inches(2.5)
width = Inches(9)
top = Inches(2)
pic = slide.shapes.add_picture(chart_fname, left, top, width=width, height=height)

# 테이블 추가
left = Inches(-1)
height = Inches(3)
width = Inches(12)
top = Inches(4)
pic = slide.shapes.add_picture(table_fname, left, top, width=width, height=height)
cursor_sp = slide.shapes[0]._element
cursor_sp.addprevious(pic._element)

#################################
## 보고서 저장
#################################
ppt_fname = os.path.join("res", "stock_report", 'stock_report.pptx')
prs.save(ppt_fname)


