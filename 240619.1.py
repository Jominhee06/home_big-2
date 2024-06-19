from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해
prs = Presentation() # 파워포인트 객체 선언

title_slide_layout = prs.slide_layouts[0] # 0 : 제목 슬라이드에 해당
slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가

# 제목 - 제목에 값넣기
title = slide.placeholders[0] # 제목
title.text = "Hello, World!" # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목 상자는 placeholders[0], 부제목 상자는 [1]
subtitle.text = "python-pptx was here!" 

# 저장
prs.save('test.pptx')

bullet_slide_layout = prs.slide_layouts[1] # 1 : 제목과 내용 슬라이드
slide = prs.slides.add_slide(bullet_slide_layout) # 기존에 있던 슬라이드에 추가

# 제목
title_shape = slide.placeholders[0]
title_shape.text = 'Adding a Bullet Slide'

# 내용
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

# 단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1 # 1 : 들여쓰기 레벨

# 단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2 # 2 : 들여쓰기 레벨

prs.save('test.pptx')

img_path = 'D:/Joo.min.hee/IOT강의/Source.1/powerpoint_handling/test_1.jpg'

blank_slide_layout = prs.slide_layouts[6] # 6 : 제목/내용이 없는 '빈' 슬라이드
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
width = height = Inches(1)
# width, height가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

left = Inches(3)
width = Inches(5.5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

prs.save('test.pptx')

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

title_shape = slide.placeholders[0] 
title_shape.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# 열 너비 설정
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# 열 제목 작성
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# 본문 셀 작성
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

prs.save('test.pptx')



