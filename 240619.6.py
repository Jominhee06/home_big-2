import datetime
from pptx import Presentation # 파워포인트 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해
import pandas as pd # 데이터프레임 사용을 위해 pandas 라이브러리 불러오기
import os
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

#################################
## 파워포인트 객체 선언
#################################
today = datetime.datetime.today().strftime('%Y%m%d')
prs = Presentation() # 파워포인트 객체 선언

#################################
## 제목 슬라이드 추가
#################################
title_slide_layout = prs.slide_layouts[0] # 0 : 제목 슬라이드에 해당
slide = prs.slides.add_slide(title_slide_layout) # 제목 슬라이드를 파워포인트 객체에 추가

# 제목 - 제목에 값넣기
title = slide.shapes.title # 제목
title.text = "주식 보고서" # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목 상자는 placeholders[0], 부제목 상자는 [1]
subtitle.text = "보고서 작성일 : {date}".format(date=today)

#################################
## 차트 및 테이블 슬라이드 추가
#################################
# 예시 변수 (실제 데이터를 사용하기 전에 정의 필요)
company = "삼성전자" # 회사 이름
df = pd.DataFrame({'close': [50000]}) # 데이터프레임 예시
chart_fname = "D:/Joo.min.hee/IOT강의/Source.1/res/stock_report/삼성전자_chart.png" # 차트 이미지 파일 경로
table_fname = "D:/Joo.min.hee/IOT강의/Source.1/res/stock_report/삼성전자_table.png" # 테이블 이미지 파일 경로

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

shapes = slide.shapes
shapes.title.text = '{company}, {close}원에 거래 마감'.format(company=company, close=df.iloc[0]['close'])
print(shapes.title.text)

# 차트 추가
left = Inches(0.5)
height = Inches(2.5)
width = Inches(9)
top = Inches(2)
# width, height가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(chart_fname, left, top, width=width, height=height)

# 테이블 추가
left = Inches(-1)
height = Inches(3)
width = Inches(12)
top = Inches(4)

pic = slide.shapes.add_picture(table_fname, left, top, width=width, height=height)
cursor_sp = slide.shapes[0]._element
cursor_sp.addprevious(pic._element) # 해당 요소를 뒤로 보냅니다.

#################################
## 보고서 저장
#################################
ppt_fname = os.path.join("res\\stock_report", 'stock_report.pptx')
prs.save(ppt_fname)

#################################
## 이메일로 보고서 보내기
#################################
# 이메일 설정
smtp_server = 'smtp.naver.com' # SMTP 서버 주소 (예: Gmail)
smtp_port = 587 # SMTP 포트 (예: Gmail의 경우 587)
sender_email = 'jominhi0617@naver.com' # 보내는 사람 이메일 주소
sender_password = 'kimsdf@@2560' # 보내는 사람 이메일 비밀번호
receiver_email = 'jominhi0617@naver.com' # 받는 사람 이메일 주소

# 이메일 메시지 작성
msg = MIMEMultipart()
msg['From'] = sender_email
msg['To'] = receiver_email
msg['Subject'] = '주식 보고서'

body = '첨부된 파일은 주식 보고서입니다.'
msg.attach(MIMEText(body, 'plain'))

# 파일 첨부
filename = 'stock_report.pptx'
attachment_path = os.path.join('res/stock_report', filename)
attachment = open(attachment_path, 'rb')
part = MIMEBase('application', 'octet-stream')
part.set_payload(attachment.read())
encoders.encode_base64(part)
part.add_header('Content-Disposition', f'attachment; filename={filename}')

msg.attach(part)

# 이메일 보내기
try:
    server = smtplib.SMTP(smtp_server, smtp_port)
    server.starttls()
    server.login(sender_email, sender_password)
    text = msg.as_string()
    server.sendmail(sender_email, receiver_email, text)
    print('이메일이 성공적으로 보내졌습니다.')
except Exception as e:
    print(f'이메일을 보내는 중 오류가 발생했습니다: {e}')
finally:
    server.quit()
    attachment.close()
